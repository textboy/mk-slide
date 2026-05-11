#!/usr/bin/env python3
"""
Extract all content from a PowerPoint file (.pptx) for HTML slide generation.

Outputs a structured JSON file with slides, text, images, tables, and speaker notes.
Images are extracted to an assets/ directory alongside the JSON output.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]
    python extract-pptx.py presentation.pptx ./extracted

Requires: pip install python-pptx Pillow

Output format: extracted-slides.json
    [
        {
            "number": 1,
            "title": "Slide Title",
            "layout": "Title Slide",
            "content": [
                {"type": "text", "content": "...", "level": 0, "bold": false, "font_size": 24},
                {"type": "table", "headers": [...], "rows": [...]},
                {"type": "group", "items": [...]}
            ],
            "images": [
                {"path": "assets/slide1_img1.png", "width": 800, "height": 600, "alt": ""}
            ],
            "notes": "Speaker notes text",
            "background": {"type": "solid", "color": "#FFFFFF"} | null
        }
    ]
"""

import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


def emu_to_px(emu_val):
    """Convert EMU (English Metric Units) to approximate pixels at 96 DPI."""
    if emu_val is None:
        return None
    return round(emu_val / 914400 * 96)


def color_to_hex(color):
    """Convert a pptx color to hex string, or return None."""
    try:
        if color and color.rgb:
            return f"#{color.rgb}"
    except (AttributeError, TypeError):
        pass
    return None


def extract_text_runs(paragraph):
    """Extract rich text from a paragraph, preserving formatting."""
    runs = []
    for run in paragraph.runs:
        run_data = {
            "text": run.text,
            "bold": run.font.bold or False,
            "italic": run.font.italic or False,
            "font_size": run.font.size.pt if run.font.size else None,
            "color": color_to_hex(run.font.color) if run.font.color else None,
        }
        runs.append(run_data)
    return runs


def extract_table(shape):
    """Extract table data from a table shape."""
    table = shape.table
    headers = []
    rows = []

    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip() for cell in row.cells]
        if row_idx == 0:
            headers = cells
        else:
            rows.append(cells)

    return {
        "type": "table",
        "headers": headers,
        "rows": rows,
    }


def extract_background(slide):
    """Attempt to extract slide background color."""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            if hasattr(fill, 'fore_color') and fill.fore_color:
                return {
                    "type": "solid",
                    "color": color_to_hex(fill.fore_color),
                }
    except Exception:
        pass
    return None


def process_image(image_path, max_dim=1600):
    """Resize oversized images to reduce output size. Requires Pillow."""
    if not HAS_PIL:
        return

    try:
        img = Image.open(image_path)
        if max(img.size) > max_dim:
            img.thumbnail((max_dim, max_dim), Image.LANCZOS)
            # Preserve format
            fmt = img.format or 'PNG'
            save_kwargs = {}
            if fmt.upper() == 'JPEG':
                save_kwargs['quality'] = 85
            img.save(image_path, fmt, **save_kwargs)
    except Exception:
        pass  # Non-fatal: keep original image


def extract_pptx(file_path, output_dir="."):
    """
    Extract all content from a PowerPoint file.

    Args:
        file_path: Path to .pptx file
        output_dir: Directory to write output (JSON + assets/)

    Returns:
        List of slide data dictionaries
    """
    file_path = Path(file_path)
    output_dir = Path(output_dir)

    if not file_path.exists():
        print(f"Error: File not found: {file_path}")
        sys.exit(1)

    if not file_path.suffix.lower() == '.pptx':
        print(f"Warning: File does not have .pptx extension: {file_path}")

    prs = Presentation(str(file_path))
    slides_data = []

    # Create assets directory
    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    # Extract slide dimensions for aspect ratio reference
    slide_width = emu_to_px(prs.slide_width)
    slide_height = emu_to_px(prs.slide_height)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "layout": "",
            "content": [],
            "images": [],
            "notes": "",
            "background": extract_background(slide),
        }

        # Extract layout name (useful for determining slide type)
        try:
            slide_data["layout"] = slide.slide_layout.name
        except Exception:
            pass

        for shape in slide.shapes:
            # --- Text content ---
            if shape.has_text_frame:
                is_title = (shape == slide.shapes.title) if slide.shapes.title else False

                if is_title:
                    slide_data["title"] = shape.text.strip()
                else:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue

                        content_item = {
                            "type": "text",
                            "content": text,
                            "level": para.level or 0,
                            "alignment": str(para.alignment) if para.alignment else None,
                            "runs": extract_text_runs(para),
                        }
                        slide_data["content"].append(content_item)

            # --- Tables ---
            if shape.has_table:
                slide_data["content"].append(extract_table(shape))

            # --- Images ---
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext or 'png'
                    img_index = len(slide_data["images"]) + 1
                    image_name = f"slide{slide_num + 1}_img{img_index}.{image_ext}"
                    image_path = assets_dir / image_name

                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    # Auto-resize oversized images
                    process_image(image_path)

                    slide_data["images"].append({
                        "path": f"assets/{image_name}",
                        "width": emu_to_px(shape.width),
                        "height": emu_to_px(shape.height),
                        "left": emu_to_px(shape.left),
                        "top": emu_to_px(shape.top),
                        "alt": shape.name or "",
                    })
                except Exception as e:
                    print(f"  Warning: Could not extract image from slide {slide_num + 1}: {e}")

            # --- Grouped shapes ---
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_items = []
                try:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            for para in child.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    group_items.append({
                                        "type": "text",
                                        "content": text,
                                        "runs": extract_text_runs(para),
                                    })
                except Exception:
                    pass

                if group_items:
                    slide_data["content"].append({
                        "type": "group",
                        "items": group_items,
                    })

        # --- Speaker notes ---
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_data["notes"] = notes_text
        except Exception:
            pass

        slides_data.append(slide_data)

    return slides_data, {"width": slide_width, "height": slide_height}


def print_summary(slides_data, dimensions):
    """Print a human-readable summary of extracted content."""
    print(f"\nSlide dimensions: {dimensions['width']}x{dimensions['height']}px")
    print(f"Total slides: {len(slides_data)}\n")

    for s in slides_data:
        title = s["title"] or "(no title)"
        img_count = len(s["images"])
        text_count = len([c for c in s["content"] if c["type"] == "text"])
        table_count = len([c for c in s["content"] if c["type"] == "table"])
        layout = s["layout"]
        notes = " [has notes]" if s["notes"] else ""

        parts = []
        if text_count:
            parts.append(f"{text_count} text")
        if img_count:
            parts.append(f"{img_count} img")
        if table_count:
            parts.append(f"{table_count} table")

        detail = ", ".join(parts) if parts else "empty"
        print(f"  Slide {s['number']:2d}: {title[:50]:<50s}  ({detail}){notes}")
        if layout:
            print(f"           Layout: {layout}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    print(f"Extracting: {input_file}")
    slides, dimensions = extract_pptx(input_file, output_dir)

    # Write extracted data as JSON
    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump({
            "dimensions": dimensions,
            "slides": slides,
        }, f, indent=2, ensure_ascii=False)

    print(f"Output: {output_path}")
    print_summary(slides, dimensions)
