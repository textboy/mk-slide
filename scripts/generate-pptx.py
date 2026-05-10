#!/usr/bin/env python3
"""
generate-pptx.py — MK Slide PPTX converter.

Extracts full visual tree from each HTML slide via Playwright:
  - Text nodes → text boxes with exact font, size, color, style, alignment
  - Containers with backgrounds → rectangles with fill color
  - Borders → shape strokes with color and width
  - Border-radius → rounded rectangles
  - Text-shadow → glow effects
  - Background color → slide fill
  - Styled inline spans (accent colors) → inline text runs
  - Outline text (-webkit-text-stroke) → text outline effect

Usage:
    python scripts/generate-pptx.py test/deck.html --output test/deck.pptx
"""

import argparse
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = 1920
SLIDE_H = 1080
PPTX_W = Inches(13.333)
PPTX_H = Inches(7.5)


def scale_x(x):
    return int(max(0, min(x, SLIDE_W)) / SLIDE_W * PPTX_W)


def scale_y(y):
    return int(max(0, min(y, SLIDE_H)) / SLIDE_H * PPTX_H)


def scale_w(w):
    return int(max(0, min(w, SLIDE_W)) / SLIDE_W * PPTX_W)


def scale_h(h):
    return int(max(0, min(h, SLIDE_H)) / SLIDE_H * PPTX_H)


def hex_to_rgb(color_str):
    if not color_str or color_str in ('transparent', 'none', 'rgba(0, 0, 0, 0)'):
        return None
    m = re.search(r'#([0-9a-fA-F]{2})([0-9a-fA-F]{2})([0-9a-fA-F]{2})', color_str)
    if m:
        return RGBColor(int(m.group(1), 16), int(m.group(2), 16), int(m.group(3), 16))
    m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*[\d.]+)?\)', color_str)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None


FONT_MAP = {
    'orbitron': 'Orbitron', 'space mono': 'Space Mono',
    'inter': 'Inter', 'archivo': 'Archivo',
    'noto sans sc': 'Microsoft YaHei', 'noto serif sc': 'SimSun',
    'system-ui': 'Calibri', 'sans-serif': 'Calibri',
    'serif': 'Times New Roman', 'monospace': 'Consolas',
}


def map_font(family):
    if not family:
        return 'Calibri'
    fl = family.lower()
    for key, val in FONT_MAP.items():
        if key in fl:
            return val
    for name in family.split(','):
        n = name.strip().strip("'\"").lower()
        if n not in ('sans-serif', 'serif', 'monospace', 'system-ui'):
            return name.strip().strip("'\"")
    return 'Calibri'


# ── JS: extract all visible elements with full computed styles ──
EXTRACT_JS = """(function(){
var slide=document.querySelectorAll('.slide')[SIDX];if(!slide)return{elements:[],bg:null};
var bg=getComputedStyle(slide).backgroundColor;
var bgNorm=(bg||'').replace(/\\s/g,'');
if(!bgNorm||bgNorm==='rgba(0,0,0,0)'||bgNorm==='transparent')bg=getComputedStyle(document.body).backgroundColor;
function buildText(n){var t='';var nodes=n.childNodes;
for(var i=0;i<nodes.length;i++){var nd=nodes[i];
if(nd.nodeType===3){t+=nd.textContent;}else if(nd.nodeType===1){var tag=nd.tagName.toLowerCase();
if(tag==='br')t+=' ';else if(tag==='span'||tag==='b'||tag==='strong'||tag==='em'||tag==='i'||tag==='u'){t+=nd.innerText||nd.textContent||'';}
else t+=buildText(nd);}}
return t.replace(/\\s+/g,' ').trim();}
// Select all visible elements including decorative text
var sel='h1,h2,h3,h4,p,li,span,b,strong,em,div,.section-label,.content-heading,.body-text,.subtitle,.title-huge,.close-title,.close-sub,.stat-block,.feature-row,.col,.pillar-card,.maturity-row,.roadmap-row,.house-roof,.house-floor,.floor-item,.quote-bar,.red-block,.hud-info,.cursor';
var all=slide.querySelectorAll(sel);var els=[];
for(var i=0;i<all.length;i++){
var n=all[i];var cs=getComputedStyle(n);var r=n.getBoundingClientRect();
if(r.width<3||r.height<3)continue;if(parseFloat(cs.opacity)<0.1)continue;if(r.left>2000||r.top>1200)continue;
if(n.closest('.progress-bar,.page-counter,.glow-orb'))continue;
var tag=n.tagName.toLowerCase();
var txt=buildText(n);
// Keep outline text (webkit-text-stroke) - render it as outline in PPTX
var sw=parseFloat(cs.webkitTextStrokeWidth)||0;var ol=sw>0;
var hasBg=cs.backgroundColor&&!cs.backgroundColor.match(/rgba\\(0,\\s*0,\\s*0,\\s*0\\)|transparent/);
var hasB=parseFloat(cs.borderTopWidth)>0&&cs.borderTopStyle!=='none';
var isTextTag=(tag==='h1'||tag==='h2'||tag==='h3'||tag==='h4'||tag==='p'||tag==='li'||tag==='span'||tag==='b'||tag==='strong'||tag==='em');
if(tag==='div'&&!hasBg&&(!hasB||parseFloat(cs.borderTopWidth)<=1.5)){
var dc=n.className&&(typeof n.className==='string'?n.className:n.className.baseVal)||'';
if(!/(^|\s)(number|label|section-label|hud-info|cursor|close-title|close-sub|quote-bar|feature-row|pillar-card|floor-item|house-roof|house-floor|col)\s/.test(dc+' '))continue;}
if(isTextTag&&!txt&&!ol)continue;
// For heading outlines, keep them as separate elements so they render with outline effect
// Dont skip outline spans inside headings - let them render as separate outlined text
var strokeColor='';
if(ol){strokeColor=cs.webkitTextStrokeColor||cs.color;}
var sh=cs.textShadow;var glow=sh&&sh!=='none'&&sh.indexOf('0px 0px 0px')===-1;
els.push({tag:tag,text:txt.slice(0,800),x:r.left,y:r.top,w:r.width,h:r.height,bg:hasBg?cs.backgroundColor:'',borderColor:cs.borderTopColor,borderWidth:parseFloat(cs.borderTopWidth)||0,borderRadius:parseFloat(cs.borderTopLeftRadius)||0,opacity:parseFloat(cs.opacity)||1,isOutline:ol,strokeColor:strokeColor,hasGlow:glow,paddingTop:parseFloat(cs.paddingTop)||0,paddingRight:parseFloat(cs.paddingRight)||0,paddingBottom:parseFloat(cs.paddingBottom)||0,paddingLeft:parseFloat(cs.paddingLeft)||0,fontFamily:cs.fontFamily,fontSize:parseFloat(cs.fontSize)||0,fontWeight:parseInt(cs.fontWeight)||400,fontStyle:cs.fontStyle||'normal',color:cs.color,textAlign:cs.textAlign||'left',textTransform:cs.textTransform||'none'});
}
// Keep outline+h1+h2 first, then sort by priority then text length
var pri={'h1':7,'h2':6,'h3':5,'h4':4,'p':3,'li':3,'span':2,'div':1,'b':1};
var byPos={};
for(var i=0;i<els.length;i++){var e=els[i];if(!e.text&&!e.bg)continue;
var k=Math.round(e.x/8)+','+Math.round(e.y/8);var p=byPos[k];var pc=pri[e.tag]||0;var pp=p?(pri[p.tag]||0):-1;
if(!p||pc>pp||(pc===pp&&e.text.length>p.text.length))byPos[k]=e;}
var items=Object.values(byPos);
items.sort(function(a,b){return b.text.length-a.text.length;});
var keep=[];
for(var i=0;i<items.length;i++){var a=items[i];
// Outline and glow elements are always kept (special visual styling would be lost)
if(a.isOutline||a.hasGlow){keep.push(a);continue;}
var dup=false;
for(var j=0;j<keep.length;j++){var b=keep[j];
if(Math.abs(a.x-b.x)<5&&Math.abs(a.y-b.y)<5){
if(b.isOutline||b.hasGlow)continue; // outline/glow are different renderings
if(b.text.indexOf(a.text)>=0&&a.text.length<b.text.length){dup=true;break;}
if(a.text.indexOf(b.text)>=0&&b.text.length<a.text.length){dup=true;break;}}}
if(!dup)keep.push(a);}
keep.sort(function(a,b){return Math.abs(a.y-b.y)<5?a.x-b.x:a.y-b.y;});
// Remove parent elements whose text is entirely covered by styled child elements
for(var i=keep.length-1;i>=0;i--){var a=keep[i];
if(a.isOutline||a.hasGlow||a.text.length===0)continue;
// Collect children in vertical bounds, sorted by text length descending
// (longest first prevents short text like "%" from breaking longer matches)
var children=[];var remaining=a.text;
for(var j=0;j<keep.length;j++){if(i===j)continue;
var b=keep[j];
if(b.y>=a.y-30&&b.y+b.h<=a.y+a.h+30&&remaining.indexOf(b.text)>=0){
children.push(b.text);}}
children.sort(function(x,y){return y.length-x.length;});
for(var ci=0;ci<children.length;ci++){remaining=remaining.replace(children[ci],'');}
remaining=remaining.replace(/\s+/g,'').trim();
if(remaining.length===0||(a.bg&&remaining.length<30)){a.text='';a.fontSize=0;}}
return{elements:keep,bg:bg};})()
"""


def extract_slide_data(page, slide_index):
    return page.evaluate(EXTRACT_JS.replace('SIDX', str(slide_index)))


def build_slide_from_elements(slide, elements, bg_color, slide_w, slide_h):
    bg_rgb = hex_to_rgb(bg_color)
    if bg_rgb:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        s.fill.solid()
        s.fill.fore_color.rgb = bg_rgb
        s.line.fill.background()

    for el in elements:
        x, y = scale_x(el['x']), scale_y(el['y'])
        w, h = scale_w(el['w']), scale_h(el['h'])
        if w <= 0 or h <= 0:
            continue

        # Clamp shapes to slide boundary — no component may exceed the deck edge
        if x + w > PPTX_W:
            w = max(0, PPTX_W - x)
        if y + h > PPTX_H:
            h = max(0, PPTX_H - y)
        if w <= 0 or h <= 0:
            continue

        has_bg = el['bg'] and el['bg'] not in ('transparent', 'rgba(0, 0, 0, 0)') and el['bg'] != ''
        has_border = el['borderWidth'] > 0 and el['borderColor'] not in ('transparent', 'rgba(0, 0, 0, 0)')
        has_text = el['text'] and el['fontSize'] > 0
        is_rounded = el['borderRadius'] > 4

        # Render shape if element has background or border
        if has_bg or has_border:
            stype = MSO_SHAPE.ROUNDED_RECTANGLE if is_rounded else MSO_SHAPE.RECTANGLE
            s = slide.shapes.add_shape(stype, x, y, w, h)
            if has_bg:
                frgb = hex_to_rgb(el['bg'])
                if frgb:
                    s.fill.solid()
                    s.fill.fore_color.rgb = frgb
            else:
                s.fill.background()
            if has_border:
                brgb = hex_to_rgb(el['borderColor'])
                if brgb:
                    s.line.color.rgb = brgb
                s.line.width = Pt(max(0.5, el['borderWidth'] * 0.5))
            else:
                s.line.fill.background()

        # Render text if present
        if has_text:
            txt = el['text']
            if el['textTransform'] == 'uppercase':
                txt = txt.upper()

            tx, ty, tw, th = x, y, w, h
            padt = scale_h(el['paddingTop']) if el['paddingTop'] else 0
            padl = scale_w(el['paddingLeft']) if el['paddingLeft'] else 0
            padb = scale_h(el['paddingBottom']) if el['paddingBottom'] else 0
            padr = scale_w(el['paddingRight']) if el['paddingRight'] else 0
            if padt or padl or padb or padr:
                tx += padl
                ty += padt
                tw -= (padl + padr)
                th -= (padt + padb)

            tx = max(0, tx)
            ty = max(0, ty)
            tw = max(10, min(tw, slide_w - tx))
            th = max(6, min(th, slide_h - ty))

            font_size = max(4, min(60, el['fontSize'] * 0.65))
            bold = el['fontWeight'] >= 600
            italic = el['fontStyle'] == 'italic'
            font_color = hex_to_rgb(el['color']) or RGBColor(200, 204, 224)
            font_name = map_font(el['fontFamily'])

            align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                        'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY}
            alignment = align_map.get(el['textAlign'], PP_ALIGN.LEFT)

            txbox = slide.shapes.add_textbox(tx, ty, tw, th)
            tf = txbox.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(3)
            tf.margin_right = Pt(3)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)

            # Build paragraph with inline runs (styled spans) or plain text
            p = tf.paragraphs[0]
            r = None
            inline_runs = el.get('inlineRuns', [])

            if inline_runs:
                # Strategy: build text by interleaving plain segments and styled runs.
                # Set p.text once (first segment), then use p.add_run() for all subsequent.
                remaining = txt
                first = True
                for ir in inline_runs:
                    ir_text = ir['text']
                    idx = remaining.find(ir_text)
                    if idx >= 0:
                        # Plain text before this run (preserve spaces)
                        prefix = remaining[:idx]
                        if prefix:
                            if first:
                                p.text = prefix
                                p.font.size = Pt(font_size)
                                p.font.color.rgb = font_color
                                p.font.name = font_name
                                first = False
                            else:
                                rr = p.add_run()
                                rr.text = prefix
                                rr.font.size = Pt(font_size)
                                rr.font.color.rgb = font_color
                                rr.font.name = font_name

                        # The styled run
                        ir_color = hex_to_rgb(ir['color']) or RGBColor(61, 90, 254)
                        ir_fs = max(8, min(60, (ir['fontSize'] or el['fontSize']) * 0.65))
                        ir_bold = ir['bold'] or bold

                        if first:
                            p.text = ir_text
                            p.font.size = Pt(ir_fs)
                            p.font.bold = ir_bold
                            p.font.color.rgb = ir_color
                            p.font.name = font_name
                            first = False
                        else:
                            rr = p.add_run()
                            rr.text = ir_text
                            rr.font.size = Pt(ir_fs)
                            rr.font.bold = ir_bold
                            rr.font.color.rgb = ir_color
                            rr.font.name = font_name

                        remaining = remaining[idx + len(ir_text):]

                # Remaining text after last run
                if remaining:
                    if first:
                        p.text = remaining
                        p.font.size = Pt(font_size)
                        p.font.color.rgb = font_color
                        p.font.name = font_name
                    else:
                        rr = p.add_run()
                        rr.text = remaining
                        rr.font.size = Pt(font_size)
                        rr.font.color.rgb = font_color
                        rr.font.name = font_name
            else:
                p.text = txt
                p.font.size = Pt(font_size)
                p.font.bold = bold
                p.font.italic = italic
                p.font.color.rgb = font_color
                p.font.name = font_name

            p.alignment = alignment

            # Outline text effect (-webkit-text-stroke)
            if el.get('isOutline'):
                try:
                    from pptx.oxml.ns import qn
                    p_elem = txbox.text_frame.paragraphs[0]._p
                    pPr = p_elem.get_or_add_pPr()
                    ln = pPr.makeelement(qn('a:ln'), {'w': '25400'})
                    sf = ln.makeelement(qn('a:solidFill'), {})
                    stroke_str = el.get('strokeColor', '')
                    if stroke_str:
                        m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', stroke_str)
                        if m:
                            val_str = f'{int(m.group(1)):02X}{int(m.group(2)):02X}{int(m.group(3)):02X}'
                        else:
                            val_str = 'C8CCE0'
                    else:
                        val_str = 'C8CCE0'
                    sc = sf.makeelement(qn('a:srgbClr'), {'val': val_str})
                    sf.append(sc)
                    ln.append(sf)
                    pPr.append(ln)
                    # Make text fill transparent so only the outline stroke is visible
                    r_elem = p_elem.find(qn('a:r'))
                    if r_elem is not None and hex_to_rgb(el.get('color', '')) is None:
                        rPr = r_elem.get_or_add_rPr()
                        for existing in rPr.findall(qn('a:solidFill')):
                            rPr.remove(existing)
                        noFill = rPr.makeelement(qn('a:noFill'), {})
                        rPr.append(noFill)
                except Exception:
                    pass

            # Glow effect for text-shadow
            if el.get('hasGlow') and bold:
                try:
                    from pptx.oxml.ns import qn
                    p_elem = txbox.text_frame.paragraphs[0]._p
                    pPr = p_elem.get_or_add_pPr()
                    glow = pPr.makeelement(qn('a:glow'), {'rad': '60000'})
                    sc = glow.makeelement(qn('a:srgbClr'), {'val': '3d5afe'})
                    glow.append(sc)
                    pPr.append(glow)
                except Exception:
                    pass


def convert_html_to_pptx(html_path, output_path):
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("Error: playwright not installed. Run: pip install playwright && playwright install chromium")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = PPTX_W
    prs.slide_height = PPTX_H
    blank_layout = prs.slide_layouts[6]
    file_url = f"file://{os.path.abspath(html_path)}"

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        page.goto(file_url, wait_until="networkidle")

        slide_count = page.evaluate("() => document.querySelectorAll('.slide').length")
        if slide_count == 0:
            slide_count = 1

        print(f"   Slides: {slide_count}")
        print(f"   Extracting visual tree...")

        for i in range(slide_count):
            if slide_count > 1:
                page.evaluate("""
                    (function(){
                        var go=window.go;if(go){go(""" + str(i) + """);return;}
                        var ss=document.querySelectorAll('.slide');
                        ss.forEach(function(s){s.classList.remove('active');});
                        if(ss[""" + str(i) + """]) ss[""" + str(i) + """].classList.add('active');
                    })()
                """)
                page.wait_for_timeout(600)

            data = extract_slide_data(page, i)
            elements = data.get('elements', [])
            bg = data.get('bg', '')

            print(f"      Slide {i+1}: {len(elements)} visual elements")

            slide = prs.slides.add_slide(blank_layout)
            build_slide_from_elements(slide, elements, bg, PPTX_W, PPTX_H)

        browser.close()

    prs.save(output_path)
    print(f"OK: {os.path.basename(html_path)} -> {os.path.basename(output_path)}")
    print(f"    Slides: {slide_count}")
    print(f"    Size:   {os.path.getsize(output_path) / 1024:.1f} KB")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert HTML to native PowerPoint shapes")
    parser.add_argument("input", help=".html file")
    parser.add_argument("--output", "-o", help="Output .pptx file")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: File not found: {args.input}")
        sys.exit(1)

    output = args.output
    if not output:
        base = os.path.splitext(os.path.basename(args.input))[0]
        output = os.path.abspath(base + '.pptx')
    elif not output.startswith('/'):
        output = os.path.abspath(output)

    convert_html_to_pptx(os.path.abspath(args.input), output)
